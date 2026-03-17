#!/usr/bin/env python3
"""
Build a presentation from slide images and upload to Google Slides.

Usage:
  python3 build_deck.py --title "My Deck" --slides s1.png s2.png --account email@example.com --upload
  python3 build_deck.py --title "My Deck" --slides s1.png --output local.pptx  # local only
"""

import argparse
import json
import os
import subprocess
import sys
import tempfile

import requests
from pptx import Presentation
from pptx.util import Inches, Emu


# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def get_access_token(account_email: str) -> str:
    """Get OAuth access token via gog token export."""
    # Export token
    tmp = tempfile.mktemp(suffix=".json")
    env = os.environ.copy()
    env_file = os.path.join(os.path.dirname(__file__), "../../../.env.gog")
    if os.path.exists(env_file):
        with open(env_file) as f:
            for line in f:
                line = line.strip()
                if "=" in line and not line.startswith("#"):
                    k, v = line.split("=", 1)
                    env[k] = v

    result = subprocess.run(
        ["gog", "auth", "tokens", "export", account_email, "--out", tmp, "--overwrite"],
        capture_output=True, text=True, env=env
    )
    if result.returncode != 0:
        print(json.dumps({"success": False, "error": f"gog export failed: {result.stderr}"}))
        sys.exit(1)

    with open(tmp) as f:
        token_data = json.load(f)
    os.unlink(tmp)

    # Load client credentials
    creds_path = os.path.expanduser("~/.config/gogcli/credentials.json")
    with open(creds_path) as f:
        creds = json.load(f)

    # Exchange refresh token for access token
    resp = requests.post("https://oauth2.googleapis.com/token", data={
        "client_id": creds["client_id"],
        "client_secret": creds["client_secret"],
        "refresh_token": token_data["refresh_token"],
        "grant_type": "refresh_token",
    })

    if resp.status_code != 200:
        print(json.dumps({"success": False, "error": f"Token exchange failed: {resp.text}"}))
        sys.exit(1)

    return resp.json()["access_token"]


def build_pptx(title: str, slide_paths: list[str], output_path: str) -> str:
    """Create a .pptx with images as full-bleed slide backgrounds."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    for img_path in slide_paths:
        if not os.path.exists(img_path):
            print(f"Warning: {img_path} not found, skipping", file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add image as full-bleed background
        slide.shapes.add_picture(
            img_path,
            Emu(0), Emu(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )

    prs.save(output_path)
    return output_path


def upload_to_drive(access_token: str, pptx_path: str, title: str) -> dict:
    """Upload .pptx to Google Drive with conversion to Google Slides."""
    import io

    # Step 1: Create file metadata (request conversion)
    metadata = {
        "name": title,
        "mimeType": "application/vnd.google-apps.presentation"
    }

    # Step 2: Multipart upload with conversion
    with open(pptx_path, "rb") as f:
        file_data = f.read()

    # Use resumable upload for reliability
    # First: initiate
    init_resp = requests.post(
        "https://www.googleapis.com/upload/drive/v3/files?uploadType=resumable&convert=true",
        headers={
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json",
            "X-Upload-Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "X-Upload-Content-Length": str(len(file_data)),
        },
        json=metadata,
    )

    if init_resp.status_code not in (200, 308):
        return {"success": False, "error": f"Upload init failed ({init_resp.status_code}): {init_resp.text[:300]}"}

    upload_url = init_resp.headers.get("Location")
    if not upload_url:
        return {"success": False, "error": "No upload URL in response"}

    # Second: upload the file
    upload_resp = requests.put(
        upload_url,
        headers={
            "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "Content-Length": str(len(file_data)),
        },
        data=file_data,
    )

    if upload_resp.status_code not in (200, 201):
        return {"success": False, "error": f"Upload failed ({upload_resp.status_code}): {upload_resp.text[:300]}"}

    file_info = upload_resp.json()
    file_id = file_info.get("id")

    # Step 3: Set sharing permissions (anyone with link can view)
    requests.post(
        f"https://www.googleapis.com/drive/v3/files/{file_id}/permissions",
        headers={
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json",
        },
        json={
            "role": "reader",
            "type": "anyone",
        },
    )

    return {
        "success": True,
        "file_id": file_id,
        "url": f"https://docs.google.com/presentation/d/{file_id}/edit",
        "view_url": f"https://docs.google.com/presentation/d/{file_id}/present",
        "name": title,
    }


def main():
    parser = argparse.ArgumentParser(description="Build presentation deck from slide images")
    parser.add_argument("--title", required=True, help="Presentation title")
    parser.add_argument("--slides", nargs="+", required=True, help="Slide image paths (in order)")
    parser.add_argument("--output", help="Local .pptx output path (optional)")
    parser.add_argument("--account", help="Google account email for upload")
    parser.add_argument("--upload", action="store_true", help="Upload to Google Drive as Slides")
    parser.add_argument("--no-share", action="store_true", help="Don't set public sharing")

    args = parser.parse_args()

    # Validate slides exist
    for s in args.slides:
        if not os.path.exists(s):
            print(json.dumps({"success": False, "error": f"Slide not found: {s}"}))
            sys.exit(1)

    # Build .pptx
    output_path = args.output or tempfile.mktemp(suffix=".pptx")
    build_pptx(args.title, args.slides, output_path)
    print(f"Built .pptx: {output_path} ({os.path.getsize(output_path)} bytes)", file=sys.stderr)

    if args.upload:
        if not args.account:
            print(json.dumps({"success": False, "error": "--account required for upload"}))
            sys.exit(1)

        # Get token
        token = get_access_token(args.account)

        # Upload
        result = upload_to_drive(token, output_path, args.title)

        # Clean up temp file if we created one
        if not args.output and os.path.exists(output_path):
            os.unlink(output_path)

        print(json.dumps(result, indent=2))

        if result.get("success"):
            sys.exit(0)
        else:
            sys.exit(1)
    else:
        print(json.dumps({
            "success": True,
            "output_path": os.path.abspath(output_path),
            "size_bytes": os.path.getsize(output_path),
            "slides_count": len(args.slides),
            "note": "Local .pptx created. Use --upload --account <email> to publish to Google Slides."
        }, indent=2))


if __name__ == "__main__":
    main()
